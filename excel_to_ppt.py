from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
from pptx.enum.text import PP_ALIGN

# open the excel
files = pd.read_excel('name.xlsx', index_col=None)
# test the index
files.iloc[0, 4]
print(files)

# open a new Powerpoint
# or Presentation() //new
prs = Presentation('name.pptx')

for i in range(0, 4):
    x = 1
    y = 3.16
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    for j in range(0, 4):
        left = Inches(x)
        top = Inches(0.8)
        if j == 0:
            x += 1.5
        else:
            x += 3
        width = Inches(5)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        data = files.iloc[i, j]
        p = tf.add_paragraph()
        if j == 0:
            p.text = str(data)
        elif j == 1:
            p.text = 'project: ' + str(data)
        elif j == 2:
            p.text = 'project: ' + str(data)
        else:
            p.text = 'project: ' + str(data)

        p.font.size = Pt(24)
        print('i', i, 'j', j)

    # reserved for loop
    # for k in range(4,5):
    left = Inches(y)
    top = Inches(6.2)
    # y += 5
    width = Inches(7)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    data = files.iloc[i, 4]
    # data = files.iloc[i, k]
    p = tf.add_paragraph()
    p.text = str(data)
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    print('i', i, 'k', 4)
    prs.save('ba.pptx')
